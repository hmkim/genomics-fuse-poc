#!/usr/bin/env python3
"""Generate PowerPoint presentation for CRAM + AWS Mountpoint S3 PoC results."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# ── Color Palette ──
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)
BLUE_DARK = RGBColor(0x1B, 0x49, 0x65)
BLUE_MED  = RGBColor(0x2C, 0x7D, 0xA0)
BLUE_LIGHT= RGBColor(0xBE, 0xE9, 0xE8)
TEAL      = RGBColor(0x5F, 0xAD, 0x56)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
GRAY_MED  = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT= RGBColor(0xF0, 0xF0, 0xF0)
RED       = RGBColor(0xE0, 0x4F, 0x4F)
ORANGE    = RGBColor(0xE8, 0x8D, 0x2A)
GREEN     = RGBColor(0x2E, 0xA0, 0x43)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ──

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=WHITE, line_spacing=1.3, font_name="맑은 고딕"):
    """Add text box with multiple lines (list of (text, bold, color_override))."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_table(slide, left, top, width, rows, cols, data, col_widths=None,
              header_color=BLUE_DARK, font_size=12):
    """Add a styled table. data = list of lists (first row = header)."""
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(left), Inches(top),
                                        Inches(width), Inches(0.4 * rows))
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "맑은 고딕"
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = GRAY_DARK
                    paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return tbl_shape

def slide_number_footer(slide, num, total=""):
    tag = f"{num}" + (f" / {total}" if total else "")
    add_text_box(slide, 12.2, 7.0, 1.0, 0.4, tag,
                 font_size=10, color=GRAY_MED, alignment=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(s, NAVY)

# Accent line
add_rect(s, 0.8, 2.2, 1.0, 0.06, BLUE_LIGHT)

add_text_box(s, 0.8, 2.4, 11.0, 1.2,
             "CRAM + AWS Mountpoint S3 PoC",
             font_size=40, color=WHITE, bold=True)
add_text_box(s, 0.8, 3.5, 11.0, 0.8,
             "UK Biobank RAP의 dxfuse 아키텍처를 AWS 환경에서 재현하고\n"
             "3가지 접근법의 기능성과 성능을 비교 검증한 결과 보고",
             font_size=20, color=BLUE_LIGHT)

add_text_box(s, 0.8, 5.0, 6.0, 0.4, "2026-02-20  |  Amazon Linux 2023 (EC2, ap-northeast-2)",
             font_size=14, color=GRAY_MED)
add_text_box(s, 0.8, 5.4, 6.0, 0.4, "Account: <YOUR_ACCOUNT_ID>",
             font_size=12, color=GRAY_MED)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: 목차
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "목차", font_size=28, color=WHITE, bold=True)

toc_items = [
    "1.  배경 및 문제 정의",
    "2.  UK Biobank RAP 현재 아키텍처",
    "3.  AWS 재현 아키텍처 개요",
    "4.  3가지 접근법 비교",
    "5.  테스트 데이터 및 환경",
    "6.  기능 테스트 결과",
    "7.  성능 벤치마크 (합성 + 실전 데이터)",
    "8.  바이트 범위 정합성 검증",
    "9.  접근법 비교 종합",
    "10.  프로덕션 확장 아키텍처 제안",
    "11.  결론 및 다음 단계",
]
add_multi_text(s, 1.5, 1.3, 10, 5.5, toc_items,
               font_size=18, color=GRAY_DARK, line_spacing=1.6)
slide_number_footer(s, 2, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: 배경 및 문제 정의
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "1. 배경 및 문제 정의", font_size=28, color=WHITE, bold=True)

# Left: 현재 상황
add_rect(s, 0.6, 1.3, 5.8, 5.5, RGBColor(0xFA, 0xFC, 0xFD), BLUE_LIGHT)
add_text_box(s, 0.8, 1.4, 5.4, 0.5, "UK Biobank RAP 현재 환경",
             font_size=18, color=BLUE_DARK, bold=True)
add_multi_text(s, 0.8, 2.0, 5.4, 4.5, [
    ("dxfuse (DNAnexus)", True, NAVY),
    "• 연구자는 EID_7654321.cram 으로 파일에 접근",
    "• 실제 데이터는 internal_id_999.cram 으로 중앙 저장소에 1카피 보관",
    "• FUSE 마운트로 클라우드 스토리지를 로컬처럼 접근",
    "• samtools가 요청하면 해당 바이트 범위만 HTTP Range로 전송",
    "",
    ("dxfuse의 한계점", True, RED),
    "• 메모리 오버헤드 → REGENIE 등에서 OOM 발생",
    "• Cloud Workstation 불안정성 보고",
    "• Go 기반 단일 구현체에 대한 종속성",
], font_size=14, color=GRAY_DARK)

# Right: PoC 목적
add_rect(s, 6.9, 1.3, 5.8, 5.5, RGBColor(0xFA, 0xFC, 0xFD), BLUE_LIGHT)
add_text_box(s, 7.1, 1.4, 5.4, 0.5, "PoC 목적",
             font_size=18, color=BLUE_DARK, bold=True)
add_multi_text(s, 7.1, 2.0, 5.4, 4.5, [
    ("AWS 네이티브 서비스로 동일 아키텍처 재현", True, NAVY),
    "",
    "① dxfuse 핵심 구조를 Python FUSE로 재현",
    "   (SQLite 메타데이터, 프리페치 엔진, 바이트 범위 읽기)",
    "",
    "② Mountpoint S3 기반 대안 검토",
    "   (Rust 커널 FUSE + 로컬 캐시)",
    "",
    "③ 3가지 접근법 비교 검증",
    "   - 기능 정확성 (samtools 호환)",
    "   - 성능 (지연시간, 처리량)",
    "   - 바이트 범위 정합성 (bit-for-bit)",
    "",
    "④ 합성 + 1000 Genomes 실전 데이터 검증",
], font_size=14, color=GRAY_DARK)
slide_number_footer(s, 3, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: UK Biobank RAP 아키텍처
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "2. UK Biobank RAP 현재 아키텍처",
             font_size=28, color=WHITE, bold=True)

# Architecture diagram as text
add_rect(s, 0.6, 1.3, 12.1, 5.6, GRAY_LIGHT)
add_multi_text(s, 0.8, 1.4, 11.7, 5.4, [
    ("연구자 워크스테이션", True, NAVY),
    "",
    ("    samtools view  EID_7654321.cram  chr22:16M-17M", False, BLUE_DARK),
    ("           │", False, GRAY_MED),
    ("           ▼", False, GRAY_MED),
    ("    /mnt/project/  ← dx mount (dxfuse)", True, BLUE_MED),
    ("           │", False, GRAY_MED),
    ("    ┌──────┴──────────────────────────────────────┐", False, GRAY_MED),
    ("    │  dxfuse (Go)                                │", True, BLUE_DARK),
    ("    │  ├─ SQLite: namespace → inode → file_id     │", False, GRAY_DARK),
    ("    │  ├─ FileHandle: presigned URL 관리           │", False, GRAY_DARK),
    ("    │  ├─ Prefetch: 순차 접근 감지 → 비동기 선행 읽기│", False, GRAY_DARK),
    ("    │  └─ HTTP Range: bytes=X-Y                   │", False, GRAY_DARK),
    ("    └──────┬──────────────────────────────────────┘", False, GRAY_MED),
    ("           │", False, GRAY_MED),
    ("           ▼", False, GRAY_MED),
    ("    DNAnexus Cloud Storage: internal_id_999.cram (1카피)", True, NAVY),
], font_size=13, color=GRAY_DARK, font_name="Consolas")

add_text_box(s, 0.8, 6.5, 11.7, 0.5,
             "핵심: 1,000명의 연구자가 동일 파일에 접근해도 S3에는 1카피만 존재 → 저장 효율 극대화",
             font_size=14, color=BLUE_DARK, bold=True)
slide_number_footer(s, 4, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: AWS 재현 아키텍처 개요
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "3. AWS 재현 아키텍처 개요",
             font_size=28, color=WHITE, bold=True)

# Three approach boxes
approaches = [
    ("Approach A", "Custom FUSE\n(eid_fuse.py)", "dxfuse 구조 충실 재현\nSQLite + 프리페치 + S3 Range", BLUE_DARK),
    ("Approach B", "Symlink +\nMountpoint S3", "커널 FUSE (Rust)\n심볼릭 링크로 EID 매핑", BLUE_MED),
    ("Approach C", "htslib S3\nPlugin", "FUSE 없이 직접 접근\nsamtools 내장 S3 지원", RGBColor(0x5F, 0x96, 0xC2)),
]

for i, (label, title, desc, color) in enumerate(approaches):
    x = 0.6 + i * 4.2
    add_rect(s, x, 1.3, 3.8, 2.5, color)
    add_text_box(s, x + 0.1, 1.4, 3.6, 0.4, label,
                 font_size=14, color=BLUE_LIGHT, bold=True)
    add_text_box(s, x + 0.1, 1.8, 3.6, 0.8, title,
                 font_size=20, color=WHITE, bold=True)
    add_text_box(s, x + 0.1, 2.7, 3.6, 0.9, desc,
                 font_size=13, color=RGBColor(0xE0, 0xEF, 0xF5))

# Arrow down
add_text_box(s, 5.5, 3.9, 2, 0.5, "▼        ▼        ▼",
             font_size=24, color=GRAY_MED, alignment=PP_ALIGN.CENTER)

# S3 bucket
add_rect(s, 0.6, 4.5, 12.1, 2.5, RGBColor(0xF0, 0xF7, 0xFA), BLUE_LIGHT)
add_text_box(s, 0.8, 4.6, 4.0, 0.5, "Amazon S3: <YOUR_BUCKET_NAME>",
             font_size=16, color=NAVY, bold=True)

add_multi_text(s, 0.8, 5.2, 3.5, 1.5, [
    ("internal/cram/", True, BLUE_DARK),
    "  internal_id_000001.cram (3.3M)",
    "  internal_id_100001.cram (16.4G)",
    "  ... (총 12 파일, 45.7 GiB)",
], font_size=12, color=GRAY_DARK, font_name="Consolas")

add_multi_text(s, 4.8, 5.2, 3.5, 1.5, [
    ("reference/GRCh38/", True, BLUE_DARK),
    "  GRCh38_chr22.fa (50M)",
    "  GRCh38_full_...hla.fa (3.1G)",
], font_size=12, color=GRAY_DARK, font_name="Consolas")

add_multi_text(s, 8.8, 5.2, 4.0, 1.5, [
    ("metadata/eid_mapping/", True, BLUE_DARK),
    "  eid_mapping.json (합성 3샘플)",
    "  eid_mapping_1kg.json (실전 3샘플)",
    "  → EID ↔ internal_id 매핑",
], font_size=12, color=GRAY_DARK, font_name="Consolas")
slide_number_footer(s, 5, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: 3가지 접근법 상세 비교
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "4. 3가지 접근법 상세 비교",
             font_size=28, color=WHITE, bold=True)

comp_data = [
    ["평가 항목", "Approach A\nCustom FUSE", "Approach B\nSymlink+Mountpoint", "Approach C\nhtslib S3"],
    ["dxfuse 유사도", "높음 ★★★", "낮음 ★", "낮음 ★"],
    ["구현 언어", "Python (fusepy)", "Rust (mount-s3)", "C (htslib)"],
    ["EID 변환 방식", "SQLite 실시간 조회", "심볼릭 링크 (정적)", "Python 래퍼"],
    ["파일시스템 투명성", "완전 (FUSE)", "완전 (symlink)", "부분적 (CLI 래퍼)"],
    ["메모리 오버헤드", "중간", "낮음", "없음"],
    ["캐싱", "프리페치 엔진", "커널 캐시 + SSD", "없음"],
    ["프로덕션 적합성", "데몬 운영 필요", "mount-s3 안정적", "stateless"],
    ["확장성", "DynamoDB로 확장", "수만 symlink 관리", "제한 없음"],
]
add_table(s, 0.6, 1.3, 12.1, 9, 4, comp_data,
          col_widths=[2.5, 3.2, 3.2, 3.2], font_size=12)
slide_number_footer(s, 6, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: 테스트 데이터
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "5. 테스트 데이터 및 환경",
             font_size=28, color=WHITE, bold=True)

# Synthetic data
add_rect(s, 0.6, 1.3, 5.8, 2.8, RGBColor(0xFA, 0xFC, 0xFD), BLUE_LIGHT)
add_text_box(s, 0.8, 1.4, 5.4, 0.4, "합성 데이터 (Phase 1)",
             font_size=16, color=BLUE_DARK, bold=True)

synth_data = [
    ["항목", "상세"],
    ["참조 게놈", "GRCh38 chr22 (50 MiB)"],
    ["샘플 수", "3개 (wgsim 합성)"],
    ["Reads/샘플", "200,000 (100K paired-end)"],
    ["CRAM 크기", "~3.3 MiB/샘플"],
    ["총 크기", "~10 MiB"],
]
add_table(s, 0.8, 1.9, 5.4, 6, 2, synth_data,
          col_widths=[1.8, 3.6], font_size=11)

# Real data
add_rect(s, 6.9, 1.3, 5.8, 2.8, RGBColor(0xFA, 0xFC, 0xFD), BLUE_LIGHT)
add_text_box(s, 7.1, 1.4, 5.4, 0.4, "1000 Genomes 실전 데이터 (Phase 2)",
             font_size=16, color=BLUE_DARK, bold=True)

real_data = [
    ["샘플", "유형", "크기"],
    ["NA06985", "High-cov WGS 30x", "16.4 GiB"],
    ["NA06986", "High-cov WGS 30x", "14.5 GiB"],
    ["HG00096", "Low-cov WGS 4x", "14.8 GiB"],
    ["참조 게놈", "GRCh38 full + decoy", "3.1 GiB"],
    ["총 크기", "", "48.8 GiB"],
]
add_table(s, 7.1, 1.9, 5.4, 6, 3, real_data,
          col_widths=[1.5, 2.2, 1.7], font_size=11)

# Environment
add_rect(s, 0.6, 4.4, 12.1, 2.7, RGBColor(0xFA, 0xFC, 0xFD), BLUE_LIGHT)
add_text_box(s, 0.8, 4.5, 5.4, 0.4, "테스트 환경",
             font_size=16, color=BLUE_DARK, bold=True)

env_data = [
    ["항목", "상세", "항목", "상세"],
    ["OS", "Amazon Linux 2023", "samtools", "1.21 (소스 빌드)"],
    ["EC2 리전", "ap-northeast-2", "mount-s3", "1.22.0"],
    ["RAM", "61 GB", "htslib", "1.21 (S3 지원)"],
    ["vCPU", "16", "fusepy", "latest"],
    ["디스크", "1 TB NVMe SSD", "boto3", "IAM Role 인증"],
]
add_table(s, 0.8, 5.0, 11.7, 6, 4, env_data,
          col_widths=[1.5, 3.6, 1.5, 5.1], font_size=11)
slide_number_footer(s, 7, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: 기능 테스트 결과
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "6. 기능 테스트 결과",
             font_size=28, color=WHITE, bold=True)

# Big pass badge
add_rect(s, 0.6, 1.3, 3.5, 1.5, GREEN)
add_text_box(s, 0.6, 1.4, 3.5, 0.6, "27 / 27",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, 0.6, 2.0, 3.5, 0.5, "ALL PASS",
             font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Breakdown
add_rect(s, 4.5, 1.3, 4.0, 1.5, RGBColor(0xF0, 0xF7, 0xFA), BLUE_LIGHT)
add_multi_text(s, 4.7, 1.4, 3.6, 1.3, [
    ("합성 데이터 (Phase 1)", True, NAVY),
    "  Local:     4 PASS",
    "  FUSE (A):  6 PASS",
    "  Symlink (B): 5 PASS",
    "  S3 Plugin (C): 2 PASS",
], font_size=12, color=GRAY_DARK, font_name="Consolas")

add_rect(s, 8.9, 1.3, 4.0, 1.5, RGBColor(0xF0, 0xF7, 0xFA), BLUE_LIGHT)
add_multi_text(s, 9.1, 1.4, 3.6, 1.3, [
    ("실전 데이터 (Phase 2)", True, NAVY),
    "  Symlink (B): 5 PASS",
    "  FUSE (A):    4 PASS",
    "  일관성 검증:  1 PASS (25,083 reads)",
], font_size=12, color=GRAY_DARK, font_name="Consolas")

# Test details table
test_details = [
    ["테스트 항목", "검증 내용", "결과"],
    ["samtools quickcheck", "CRAM 파일 무결성 (매직 바이트, EOF 마커)", "✓ PASS"],
    ["samtools view -H", "헤더 메타데이터 (@HD, @SQ, @RG) 읽기", "✓ PASS"],
    ["samtools view (영역)", "특정 게놈 영역의 reads 추출 (바이트 범위 읽기)", "✓ PASS"],
    ["samtools flagstat", "전체 파일 순차 읽기 통계", "✓ PASS"],
    ["EID 변환 검증", "EID_XXXXXXX → internal_id 올바른 매핑", "✓ PASS"],
    ["크로스 접근법 일관성", "동일 영역 쿼리 시 Approach A = B 결과 동일", "✓ PASS"],
]
add_table(s, 0.6, 3.1, 12.1, 7, 3, test_details,
          col_widths=[3.0, 6.5, 2.6], font_size=13)

add_text_box(s, 0.8, 6.5, 11.7, 0.5,
             "실전 데이터: NA06985 (30x WGS, 16.4 GiB) chr22:16M-16.1M → 25,083 reads, 두 접근법 결과 완전 일치",
             font_size=13, color=BLUE_DARK, bold=True)
slide_number_footer(s, 8, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: 성능 벤치마크 - 합성 데이터
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "7-1. 성능 벤치마크: 합성 데이터 (3.3 MiB)",
             font_size=28, color=WHITE, bold=True)

# Results table
perf_data = [
    ["테스트", "Local\n(기준선)", "Approach A\nCustom FUSE", "Approach B\nSymlink+MP", "Approach C\nS3 Plugin"],
    ["헤더 읽기", "3 ms", "35 ms", "5 ms", "841 ms"],
    ["영역 읽기\n(chr22:1-100K)", "5 ms", "92 ms", "7 ms", "864 ms"],
    ["전체 파일\nflagstat", "27 ms", "416 ms", "31 ms", "888 ms"],
]
add_table(s, 0.6, 1.3, 12.1, 4, 5, perf_data,
          col_widths=[2.4, 2.4, 2.4, 2.4, 2.5], font_size=14)

# Relative performance
add_text_box(s, 0.8, 3.6, 10, 0.5, "상대 성능 비교 (Local = 1.0x)",
             font_size=16, color=NAVY, bold=True)

rel_data = [
    ["테스트", "Local", "FUSE (A)", "Symlink (B)", "S3 Plugin (C)"],
    ["헤더 읽기", "1.0x", "11.7x", "1.7x", "280x"],
    ["영역 읽기", "1.0x", "18.4x", "1.4x", "173x"],
    ["전체 파일", "1.0x", "15.4x", "1.1x", "32.9x"],
]
add_table(s, 0.6, 4.1, 12.1, 4, 5, rel_data,
          col_widths=[2.4, 2.4, 2.4, 2.4, 2.5], font_size=14)

add_text_box(s, 0.8, 6.3, 12.0, 0.8,
             "→ Approach B가 로컬 대비 1.1~1.7x 오버헤드로 가장 빠름  |  "
             "mount-s3의 Rust 커널 FUSE + 로컬 캐시가 핵심",
             font_size=15, color=BLUE_DARK, bold=True)
slide_number_footer(s, 9, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: 성능 벤치마크 - 실전 데이터
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "7-2. 성능 벤치마크: 1000 Genomes 실전 데이터 (16.4 GiB)",
             font_size=28, color=WHITE, bold=True)

# Real data results
real_perf = [
    ["테스트", "Approach A\nCustom FUSE", "Approach B\nSymlink+Mountpoint", "Approach C\nhtslib S3 Plugin"],
    ["헤더 읽기", "80 ms", "18 ms ★", "622 ms"],
    ["영역 읽기\n(chr22:16M-16.1M\n25,083 reads)", "508 ms", "56 ms ★", "1,301 ms"],
]
add_table(s, 0.6, 1.3, 12.1, 3, 4, real_perf,
          col_widths=[3.0, 3.0, 3.0, 3.1], font_size=15)

# Comparison table
add_text_box(s, 0.8, 3.3, 10, 0.5, "합성 vs 실전 데이터 성능 비교 (Symlink 기준)",
             font_size=16, color=NAVY, bold=True)

comp_perf = [
    ["테스트", "합성 (3.3 MiB)\nB 대비 A 비율", "실전 (16.4 GiB)\nB 대비 A 비율", "변화"],
    ["헤더 읽기", "7.0x", "4.4x", "↓ 개선"],
    ["영역 읽기", "13.1x", "9.1x", "↓ 개선"],
]
add_table(s, 0.6, 3.8, 12.1, 3, 4, comp_perf,
          col_widths=[3.0, 3.0, 3.0, 3.1], font_size=14)

# Analysis boxes
add_rect(s, 0.6, 5.2, 5.8, 2.0, RGBColor(0xE8, 0xF5, 0xE9))
add_multi_text(s, 0.8, 5.3, 5.4, 1.8, [
    ("대용량 파일에서 FUSE 성능 상대적 개선", True, GREEN),
    "",
    "• 초기 연결 오버헤드가 대형 파일에서 희석",
    "• 프리페치 엔진의 순차 읽기 최적화 효과 증가",
    "• Go/Rust 재작성 시 Approach B 수준 달성 가능",
], font_size=12, color=GRAY_DARK)

add_rect(s, 6.9, 5.2, 5.8, 2.0, RGBColor(0xFD, 0xF2, 0xE4))
add_multi_text(s, 7.1, 5.3, 5.4, 1.8, [
    ("S3 Plugin이 느린 이유", True, ORANGE),
    "",
    "• 매 호출마다 새 S3 연결 수립 (connection reuse 없음)",
    "• htslib 내부 버퍼링 전략이 FUSE 캐시 대비 비효율적",
    "• 인터랙티브 사용에는 부적합하지만 배치 처리에 적합",
], font_size=12, color=GRAY_DARK)
slide_number_footer(s, 10, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11: 바이트 범위 정합성 검증
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "8. 바이트 범위 정합성 검증",
             font_size=28, color=WHITE, bold=True)

# Badge
add_rect(s, 0.6, 1.3, 3.0, 1.2, GREEN)
add_text_box(s, 0.6, 1.4, 3.0, 0.5, "18 / 18",
             font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(s, 0.6, 1.9, 3.0, 0.4, "ALL PASS",
             font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s, 4.0, 1.5, 8.0, 0.8,
             "모든 접근 경로에서 동일 오프셋의 바이트 데이터가\n"
             "bit-for-bit 일치함을 확인 (MD5 해시 비교)",
             font_size=15, color=GRAY_DARK)

# Real data table
add_text_box(s, 0.8, 2.8, 10, 0.4, "실전 데이터 검증 (NA06985, 16.4 GiB)",
             font_size=14, color=NAVY, bold=True)

byte_data = [
    ["테스트", "오프셋", "크기", "MD5 (3소스 일치)"],
    ["CRAM magic", "0", "26B", "7b6ef75f"],
    ["First 4KB", "0", "4,096B", "ab727ab3"],
    ["1MB offset", "1,048,576", "4,096B", "0942d5b2"],
    ["100MB offset", "104,857,600", "4,096B", "9a76ffd0"],
    ["1GB offset", "1,073,741,824", "8,192B", "39caa53f"],
    ["5GB offset", "5,368,709,120", "8,192B", "9f87a066"],
    ["10GB offset", "10,737,418,240", "8,192B", "b9a2525f"],
    ["15GB offset", "16,106,127,360", "8,192B", "b03fb28d"],
    ["CRAM EOF", "17,567,630,631", "38B", "a4d9dc2f"],
]
add_table(s, 0.6, 3.2, 12.1, 10, 4, byte_data,
          col_widths=[2.5, 3.3, 2.0, 4.3], font_size=12)

add_text_box(s, 0.8, 6.8, 11.7, 0.5,
             "→ 15GB 오프셋에서도 FUSE / Symlink / S3 Direct 3소스 완전 일치: "
             "대용량 CRAM 임의 슬라이스 디코딩 정확성 보장",
             font_size=13, color=BLUE_DARK, bold=True)
slide_number_footer(s, 11, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 12: 접근법 비교 종합
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "9. 접근법 비교 종합",
             font_size=28, color=WHITE, bold=True)

summary_data = [
    ["평가 항목", "Approach A\nCustom FUSE", "Approach B\nSymlink+Mountpoint", "Approach C\nhtslib S3"],
    ["dxfuse 아키텍처\n유사도", "★★★ 높음\n(SQLite, 프리페치, FH)", "★ 낮음\n(단순 symlink)", "★ 낮음\n(FUSE 없음)"],
    ["성능\n(실전 영역 읽기)", "508 ms\n(보통)", "56 ms ★\n(최고)", "1,301 ms\n(낮음)"],
    ["메모리 오버헤드", "중간\n(Python+SQLite)", "낮음 ★\n(Rust mount-s3)", "없음 ★\n(stateless)"],
    ["EID 동적 변환", "★★★ 실시간\n(SQLite lookup)", "★ 정적\n(symlink 재생성)", "★★ 실시간\n(Python 래퍼)"],
    ["도입 복잡도", "높음\n(FUSE 데몬 운영)", "낮음 ★\n(mount-s3 + 스크립트)", "낮음 ★\n(Python)"],
    ["프로덕션 권장", "PoC / 참조용", "★★★ 권장", "배치 처리용"],
]
add_table(s, 0.6, 1.3, 12.1, 7, 4, summary_data,
          col_widths=[2.5, 3.2, 3.2, 3.2], font_size=12)

# Recommendation box
add_rect(s, 0.6, 5.8, 12.1, 1.2, RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(s, 0.8, 5.9, 11.7, 0.3, "권장 사항",
             font_size=16, color=GREEN, bold=True)
add_text_box(s, 0.8, 6.3, 11.7, 0.5,
             "프로덕션에서는 Approach B (Symlink + Mountpoint S3)를 기반으로 하되, "
             "DynamoDB를 활용한 동적 EID 매핑을 결합하여 "
             "Approach A의 실시간 변환 기능과 Approach B의 성능을 모두 확보하는 하이브리드 구조를 권장",
             font_size=14, color=GRAY_DARK)
slide_number_footer(s, 12, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 13: 프로덕션 확장 아키텍처
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "10. 프로덕션 확장 아키텍처 제안",
             font_size=28, color=WHITE, bold=True)

add_text_box(s, 0.8, 1.2, 11.7, 0.5,
             "Approach B 성능 + Approach A 동적 매핑 = 하이브리드 구조",
             font_size=18, color=NAVY, bold=True)

# Flow diagram
steps = [
    ("Cognito 인증", BLUE_DARK, "User Pool 로그인 + Group 확인\n→ Identity Pool → 프로젝트별 IAM Role"),
    ("Lambda: EID 매핑 생성", BLUE_MED, "DynamoDB에서 프로젝트별\nEID → internal_id 조회"),
    ("Mountpoint S3 마운트", RGBColor(0x5F, 0x96, 0xC2), "내부 CRAM 데이터 접근\n(커널 FUSE + SSD 캐시)"),
    ("Symlink 자동 생성", TEAL, "EID_XXXXXXX.cram →\n/mnt/s3/internal_id.cram"),
    ("연구자 접근 가능", GREEN, "samtools view EID.cram\n→ 투명한 S3 바이트 범위 읽기"),
]

for i, (title, color, desc) in enumerate(steps):
    y = 1.9 + i * 1.0
    # Step number circle
    add_rect(s, 0.6, y, 0.5, 0.5, color)
    add_text_box(s, 0.6, y + 0.05, 0.5, 0.4, str(i+1),
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(s, 1.3, y, 3.5, 0.5, title,
                 font_size=15, color=color, bold=True)
    # Description
    add_text_box(s, 1.3, y + 0.35, 3.5, 0.6, desc,
                 font_size=11, color=GRAY_MED)
    # Arrow
    if i < len(steps) - 1:
        add_text_box(s, 0.7, y + 0.55, 0.4, 0.3, "↓",
                     font_size=14, color=GRAY_MED, alignment=PP_ALIGN.CENTER)

# AWS Services mapping
add_rect(s, 5.5, 1.9, 7.2, 5.2, RGBColor(0xFA, 0xFC, 0xFD), BLUE_LIGHT)
add_text_box(s, 5.7, 2.0, 6.8, 0.4, "AWS 서비스 매핑",
             font_size=16, color=NAVY, bold=True)

svc_data = [
    ["역할", "PoC 구현", "프로덕션 확장"],
    ["EID 매핑 DB", "JSON 파일", "DynamoDB (글로벌 테이블)"],
    ["세션 초기화", "수동 스크립트", "Lambda (자동)"],
    ["CRAM 스토리지", "S3 Standard", "S3 Intelligent-Tiering"],
    ["파일시스템", "mount-s3 + symlink", "Mountpoint S3 (검증됨)"],
    ["암호화", "미적용", "S3-SSE-KMS"],
    ["감사 로그", "미적용", "CloudTrail + S3 Access Log"],
    ["인증", "IAM Role", "Cognito User Pool\n+ Identity Pool"],
]
add_table(s, 5.7, 2.5, 6.8, 8, 3, svc_data,
          col_widths=[1.7, 2.3, 2.8], font_size=11)
slide_number_footer(s, 13, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 14: 결론
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)

add_rect(s, 0.6, 0.4, 12.1, 0.06, BLUE_LIGHT)
add_text_box(s, 0.8, 0.6, 11.0, 0.8, "11. 결론",
             font_size=32, color=WHITE, bold=True)

conclusions = [
    ("1", "dxfuse 아키텍처 AWS 재현 성공",
     "Custom FUSE(eid_fuse.py)로 SQLite 메타데이터 DB, EID→S3 key 해석,\n"
     "프리페치 상태머신 등 dxfuse 핵심 구조를 충실히 구현"),
    ("2", "3가지 접근법 모두 기능적으로 정확",
     "합성 17개 + 실전 10개 = 총 27개 기능 테스트 전수 통과\n"
     "합성 9개 + 실전 9개 = 총 18개 바이트 범위 검증 bit-for-bit 일치"),
    ("3", "실전 16.4 GiB WGS 데이터에서 검증 완료",
     "1000 Genomes 30x 고심도 WGS CRAM에서 25,083 reads 정확 추출\n"
     "15GB 오프셋까지 3소스 바이트 범위 정합성 확인"),
    ("4", "Approach B (Symlink+Mountpoint S3) 최적 성능",
     "실전 데이터: 헤더 18ms, 영역 쿼리 56ms (FUSE 대비 4.4~9.1x 빠름)\n"
     "Rust 기반 mount-s3 커널 FUSE + 로컬 캐시가 핵심 성능 요인"),
    ("5", "하이브리드 프로덕션 아키텍처 제안",
     "Approach B 성능 + DynamoDB 동적 매핑 결합\n"
     "dxfuse OOM/불안정성 회피 + EID 실시간 변환 + 감사 로그 지원"),
]

for i, (num, title, desc) in enumerate(conclusions):
    y = 1.6 + i * 1.15
    add_rect(s, 0.6, y, 0.5, 0.5, BLUE_MED)
    add_text_box(s, 0.6, y + 0.05, 0.5, 0.4, num,
                 font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(s, 1.3, y, 5.0, 0.5, title,
                 font_size=17, color=BLUE_LIGHT, bold=True)
    add_text_box(s, 1.3, y + 0.4, 11.0, 0.7, desc,
                 font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC))

slide_number_footer(s, 14, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 15: 다음 단계
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_rect(s, 0, 0, 13.333, 1.0, NAVY)
add_text_box(s, 0.8, 0.2, 10, 0.6, "다음 단계",
             font_size=28, color=WHITE, bold=True)

next_steps = [
    ("단기", BLUE_DARK, [
        "동시 접근 테스트: 다수 프로세스가 동시에 서로 다른 EID를 읽는 시나리오",
        "전체 게놈 flagstat: 17 GiB 파일 전체 순차 읽기 성능 측정",
        "비용 분석: S3 GET 요청 수/데이터 전송량 기반 비용 산출",
    ]),
    ("중기", BLUE_MED, [
        "Go/Rust FUSE 재작성: eid_fuse.py Python 오버헤드 제거 → Approach B 수준 성능",
        "DynamoDB 연동: 동적 EID 매핑 DB 구현 (프로젝트별 다중 매핑)",
        "Lambda 세션 초기화: 연구자 로그인 시 EID symlink 자동 생성",
    ]),
    ("장기", TEAL, [
        "Cognito 통합: 외부 연구자 인증 + 프로젝트별 접근 제어 (Group → IAM Role)",
        "S3 Intelligent-Tiering: 접근 빈도 기반 자동 비용 최적화",
        "CloudTrail 감사 로그: CRAM 파일 접근 이력 추적 및 컴플라이언스",
    ]),
]

for i, (phase, color, items) in enumerate(next_steps):
    y = 1.3 + i * 2.0
    add_rect(s, 0.6, y, 1.5, 0.5, color)
    add_text_box(s, 0.6, y + 0.05, 1.5, 0.4, phase,
                 font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text_box(s, 2.3, y + j * 0.5, 10.5, 0.5, f"•  {item}",
                     font_size=14, color=GRAY_DARK)
slide_number_footer(s, 15, 16)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 16: Q&A
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, NAVY)

add_text_box(s, 0, 2.5, 13.333, 1.5, "Q & A",
             font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(s, 5.5, 4.0, 2.3, 0.06, BLUE_LIGHT)
add_text_box(s, 0, 4.3, 13.333, 0.8,
             "CRAM + AWS Mountpoint S3 PoC 결과 보고",
             font_size=18, color=BLUE_LIGHT, alignment=PP_ALIGN.CENTER)

add_text_box(s, 0, 5.5, 13.333, 0.5,
             "감사합니다",
             font_size=24, color=GRAY_MED, alignment=PP_ALIGN.CENTER)

# ── Save ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
OUTPUT = os.path.join(PROJECT_DIR, "PRESENTATION.pptx")
prs.save(OUTPUT)
print(f"Presentation saved: {OUTPUT}")

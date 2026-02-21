#!/usr/bin/env python3
"""Update PRESENTATION.pptx with CDK Infrastructure and Cost Analysis slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

# ── Constants ──
SLIDE_W = 12191695
SLIDE_H = 6858000
HEADER_COLOR = RGBColor(0x0D, 0x1B, 0x2A)  # Dark navy
ACCENT_TEAL = RGBColor(0xBE, 0xE9, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xFA, 0xFC, 0xFD)
GREEN_ACCENT = RGBColor(0x2E, 0xA0, 0x43)
BLUE_1 = RGBColor(0x1B, 0x49, 0x65)
BLUE_2 = RGBColor(0x2C, 0x7D, 0xA0)
BLUE_3 = RGBColor(0x5F, 0x96, 0xC2)
RED_ACCENT = RGBColor(0xE0, 0x4F, 0x5F)
FONT_NAME = "맑은 고딕"

MARGIN_LEFT = Emu(731520)
CONTENT_TOP = Emu(1188720)


def add_header_bar(slide):
    """Add the dark navy header bar at top."""
    shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W, Emu(914400)  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_COLOR
    shape.line.fill.background()
    return shape


def add_title(slide, text):
    """Add white title text in the header bar."""
    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Emu(182880), Emu(9144000), Emu(548640))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME


def add_page_number(slide, num, total):
    """Add page number at bottom right."""
    txBox = slide.shapes.add_textbox(Emu(11155680), Emu(6400800), Emu(914400), Emu(365760))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_TEXT
    run.font.name = FONT_NAME


def add_textbox(slide, left, top, width, height, texts, font_size=13,
                line_spacing=1.3, color=DARK_TEXT, bold_first=False):
    """Add a textbox with multiple lines of text."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Check for bold prefix
        is_bold = bold_first and i == 0
        actual_color = color

        if text.startswith("**") and text.endswith("**"):
            text = text[2:-2]
            is_bold = True
        elif text.startswith("[BOLD]"):
            text = text[6:]
            is_bold = True
        elif text.startswith("[GREEN]"):
            text = text[7:]
            actual_color = GREEN_ACCENT
        elif text.startswith("[RED]"):
            text = text[5:]
            actual_color = RED_ACCENT
        elif text.startswith("[GRAY]"):
            text = text[6:]
            actual_color = GRAY_TEXT

        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = actual_color
        run.font.name = FONT_NAME
        if is_bold:
            run.font.bold = True

    return txBox


def add_bg_rect(slide, left, top, width, height, color):
    """Add a colored background rectangle."""
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_numbered_box(slide, left, top, number, color):
    """Add a numbered circle/square like the production architecture slide."""
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(457200), Emu(457200))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME


def build_cdk_infra_slide(slide, page_num, total):
    """Slide: CDK 인프라 구현."""
    add_header_bar(slide)
    add_title(slide, "10. CDK 인프라 구현")

    # Left panel: Stack structure
    add_bg_rect(slide, 548640, 1188720, 5303520, 5029200, LIGHT_BG)
    add_textbox(slide, 731520, 1280160, 4937760, 457200,
                ["[BOLD]배포된 스택 구조 (5개 CDK 스택)"], font_size=16, color=BLUE_1)

    add_textbox(slide, 731520, 1737360, 4937760, 3657600, [
        "[BOLD]NetworkStack",
        "  VPC 10.0.0.0/16, 2 AZ, NAT Gateway",
        "  S3 + DynamoDB Gateway Endpoint",
        "",
        "[BOLD]StorageStack",
        "  S3 (SSE-KMS) + KMS + CloudTrail",
        "",
        "[BOLD]DatabaseStack",
        "  DynamoDB EID 매핑 (PK=project_id, SK=eid)",
        "",
        "[BOLD]AuthStack",
        "  Cognito User Pool + Identity Pool",
        "  프로젝트별 IAM Role 매핑",
        "",
        "[BOLD]ComputeStack",
        "  Lambda ×3 + EC2 워크스테이션 (t3.large)",
    ], font_size=11, line_spacing=1.2)

    # Right panel: E2E verification
    add_bg_rect(slide, 6309360, 1188720, 5303520, 5029200, LIGHT_BG)
    add_textbox(slide, 6492240, 1280160, 4937760, 457200,
                ["[BOLD]End-to-End 검증 결과"], font_size=16, color=GREEN_ACCENT)

    add_textbox(slide, 6492240, 1737360, 4937760, 4114800, [
        "[GREEN]✓ DynamoDB 시딩 완료",
        "  data-seeder Lambda → 4개 프로젝트, 10개 EID",
        "",
        "[GREEN]✓ EID 해석 검증",
        "  eid-resolver: EID_1234567 → internal_id_000001",
        "",
        "[GREEN]✓ 세션 초기화 검증",
        "  session-init: 프로젝트별 CRAM+CRAI 6개 매핑",
        "",
        "[GREEN]✓ 데이터 마이그레이션",
        "  PoC → Prod 버킷: 48.8 GiB @ 273 MiB/s",
        "",
        "[GREEN]✓ EC2 워크스테이션 검증",
        "  SSM → mount-s3 → symlink → samtools 읽기 성공",
        "",
        "[BOLD]배포 리전: ap-northeast-2",
        "[BOLD]계정: <YOUR_ACCOUNT_ID>",
    ], font_size=11, line_spacing=1.2)

    add_page_number(slide, page_num, total)


def build_cost_overview_slide(slide, page_num, total):
    """Slide: 비용 분석 - 시나리오 및 가정."""
    add_header_bar(slide)
    add_title(slide, "11. 비용 분석: 시나리오 및 가정")

    # Left: Scenarios
    add_bg_rect(slide, 548640, 1188720, 5303520, 5029200, LIGHT_BG)
    add_textbox(slide, 731520, 1280160, 4937760, 457200,
                ["[BOLD]비교 시나리오"], font_size=16, color=BLUE_1)

    add_textbox(slide, 731520, 1737360, 4937760, 4114800, [
        "[BOLD]A: Mountpoint S3 + Symlink (공유)",
        "  • S3에 데이터 1카피 저장",
        "  • 100명 연구자가 Mountpoint S3로 공유 접근",
        "  • 공유 EC2 워크스테이션 1대 (24/7)",
        "  • Lambda/DynamoDB로 EID 동적 해석",
        "",
        "[BOLD]B-1: EFS 개별 복사 (전통적 방식)",
        "  • S3 소스 + 연구자별 EFS 전체 복사",
        "  • 개별 EC2 100대 (8시간/일 × 22일)",
        "  • 스토리지: S3 × 1 + EFS × 100",
        "",
        "[BOLD]B-2: EBS 개별 복사 (전통적 방식)",
        "  • S3 소스 + 연구자별 EBS gp3 전체 복사",
        "  • 개별 EC2 100대 (8시간/일 × 22일)",
        "  • 스토리지: S3 × 1 + EBS × 100",
    ], font_size=11, line_spacing=1.2)

    # Right: Pricing & Assumptions
    add_bg_rect(slide, 6309360, 1188720, 5303520, 5029200, LIGHT_BG)
    add_textbox(slide, 6492240, 1280160, 4937760, 457200,
                ["[BOLD]AWS 서울 리전 단가 (USD)"], font_size=16, color=BLUE_2)

    add_textbox(slide, 6492240, 1737360, 4937760, 4114800, [
        "[BOLD]스토리지 (GB당 월 비용)",
        "  S3 Standard:   $0.025   (기준)",
        "  EBS gp3:         $0.0912  (3.6x)",
        "  EFS Standard:  $0.36      (14.4x)",
        "",
        "[BOLD]컴퓨팅",
        "  EC2 t3.large: $0.104/시간",
        "  Lambda: $0.20/1M건 + $0.0000167/GB-초",
        "",
        "[BOLD]기타",
        "  DynamoDB RRU: $0.375/100만건",
        "  Mountpoint S3: 무료 (추가 비용 없음)",
        "  S3 VPC Endpoint 전송: 무료",
        "",
        "[BOLD]접근 패턴 가정",
        "  연구자당 일일 50회 samtools 쿼리",
        "  쿼리당 S3 GET 100건 → 월 1,100만 GET",
    ], font_size=11, line_spacing=1.2)

    add_page_number(slide, page_num, total)


def build_cost_results_slide(slide, page_num, total):
    """Slide: 비용 분석 - 결과 비교."""
    add_header_bar(slide)
    add_title(slide, "12. 비용 분석: 비교 결과")

    # Top-left: Small project
    add_bg_rect(slide, 548640, 1188720, 5303520, 2194560, LIGHT_BG)
    add_textbox(slide, 731520, 1280160, 4937760, 365760,
                ["[BOLD]소규모 프로젝트 (50 GiB, 100 연구자)"], font_size=14, color=BLUE_1)

    add_textbox(slide, 731520, 1645920, 4937760, 1554480, [
        "                    월 비용      연 비용     배율",
        "[GREEN]A: S3+Symlink    $81       $974      기준선",
        "[RED]B-1: EFS 복사     $3,632    $43,580   44.8x",
        "[RED]B-2: EBS 복사     $2,288    $27,452   28.2x",
        "",
        "→ A 대비 B-2: 연간 $26,478 절감 (96.5%)",
    ], font_size=11, line_spacing=1.3)

    # Top-right: Production project
    add_bg_rect(slide, 6309360, 1188720, 5303520, 2194560, LIGHT_BG)
    add_textbox(slide, 6492240, 1280160, 4937760, 365760,
                ["[BOLD]프로덕션 프로젝트 (15 TB, 100 연구자)"], font_size=14, color=BLUE_1)

    add_textbox(slide, 6492240, 1645920, 4937760, 1554480, [
        "                       월 비용         연 비용        배율",
        "[GREEN]A: S3+Symlink      $464        $5,567      기준선",
        "[RED]B-1: EFS 복사    $555,180   $6,662,162   1,197x",
        "[RED]B-2: EBS 복사    $142,303   $1,707,640     307x",
        "",
        "→ A 대비 B-2: 연간 $1,702,073 절감 (99.7%)",
    ], font_size=11, line_spacing=1.3)

    # Bottom: Key insight box
    add_bg_rect(slide, 548640, 3566160, 11064480, 2743200, LIGHT_BG)
    add_textbox(slide, 731520, 3657600, 10698720, 365760,
                ["[BOLD]비용 구조 분석 (프로덕션 15 TB 기준)"], font_size=14, color=BLUE_2)

    # Left diagram
    add_textbox(slide, 731520, 4023360, 4937760, 2194560, [
        "[BOLD]Scenario A: S3 + Symlink",
        "",
        "  S3 스토리지 (1카피)     $384  (82.8%)",
        "  ████████████████████████████",
        "  EC2 공유 워크스테이션     $76  (16.4%)",
        "  █████",
        "  Lambda + DynamoDB         $0.12 (0.0%)",
        "",
        "[GREEN]합계: $464/월 (연구자당 $4.64)",
    ], font_size=11, line_spacing=1.2)

    # Right diagram
    add_textbox(slide, 6492240, 4023360, 4937760, 2194560, [
        "[BOLD]Scenario B-2: EBS 개별 복사",
        "",
        "  EBS 스토리지 (100카피) $140,083  (98.4%)",
        "  ████████████████████████████████████",
        "  EC2 개별 100대           $1,830    (1.3%)",
        "  █",
        "  S3 소스                    $384    (0.3%)",
        "",
        "[RED]합계: $142,303/월 (연구자당 $1,423)",
    ], font_size=11, line_spacing=1.2)

    add_page_number(slide, page_num, total)


def build_cost_curve_slide(slide, page_num, total):
    """Slide: 비용 분석 - 연구자 수별 비용 곡선."""
    add_header_bar(slide)
    add_title(slide, "13. 비용 분석: 연구자 수별 비용 곡선")

    # Main content area
    add_bg_rect(slide, 548640, 1188720, 7223760, 5029200, LIGHT_BG)
    add_textbox(slide, 731520, 1280160, 6858000, 365760,
                ["[BOLD]연구자 수에 따른 월별 비용 비교 (15 TB 기준)"], font_size=14, color=BLUE_1)

    add_textbox(slide, 731520, 1645920, 6858000, 4297680, [
        " 연구자     A: S3+Symlink     B-1: EFS         B-2: EBS        최저",
        " ─────────────────────────────────────────────────────────",
        "    1명         $460            $5,932           $1,803          A",
        "    5명         $460           $28,124           $7,480          A",
        "   10명         $460           $55,863          $14,575          A",
        "   25명         $461          $139,082          $35,862          A",
        "   50명         $462          $277,779          $71,341          A",
        "  100명         $464          $555,174         $142,298          A",
        "  200명         $468        $1,109,965         $284,211          A",
        "  500명         $479        $2,774,336         $709,952          A",
        "",
        "[GREEN]Scenario A: 1명 → 500명으로 증가해도 $460 → $479 (거의 변동 없음)",
        "[RED]Scenario B: 연구자 수에 정비례하여 선형 증가",
        "",
        "[BOLD]손익분기점: 존재하지 않음 (연구자 1명에서도 Scenario A가 유리)",
    ], font_size=11, line_spacing=1.25)

    # Right: Key insights
    add_bg_rect(slide, 8046720, 1188720, 3566160, 5029200, LIGHT_BG)
    add_textbox(slide, 8229600, 1280160, 3200400, 365760,
                ["[BOLD]핵심 인사이트"], font_size=14, color=GREEN_ACCENT)

    add_textbox(slide, 8229600, 1737360, 3200400, 4114800, [
        "[BOLD]비용 절감율",
        "",
        "  소규모 (50 GiB):",
        "  EFS 대비  97.8%",
        "  EBS 대비  96.5%",
        "",
        "  프로덕션 (15 TB):",
        "  EFS 대비  99.9%",
        "  EBS 대비  99.7%",
        "",
        "[BOLD]연간 절감액 (15TB)",
        "",
        "  vs EFS: $6,656,595",
        "  vs EBS: $1,702,073",
        "",
        "[BOLD]연구자당 월 비용",
        "",
        "  A:          $4.64",
        "[RED]  B-1(EFS): $5,552",
        "[RED]  B-2(EBS): $1,423",
        "",
        "  → A가 300~1,200배 저렴",
    ], font_size=11, line_spacing=1.15)

    add_page_number(slide, page_num, total)


def update_toc_slide(slide, total_slides):
    """Update table of contents slide."""
    # Find and update the TOC text box
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            if "배경 및 문제 정의" in full_text and "결론" in full_text:
                # This is the TOC text box - clear and rewrite
                tf = shape.text_frame
                # Clear existing paragraphs
                for p in tf.paragraphs:
                    p.clear()

                toc_items = [
                    "1.   배경 및 문제 정의",
                    "2.   UK Biobank RAP 현재 아키텍처",
                    "3.   AWS 재현 아키텍처 개요",
                    "4.   3가지 접근법 비교",
                    "5.   테스트 데이터 및 환경",
                    "6.   기능 테스트 결과",
                    "7.   성능 벤치마크 (합성 + 실전 데이터)",
                    "8.   바이트 범위 정합성 검증",
                    "9.   접근법 비교 종합",
                    "10.  프로덕션 확장 아키텍처 제안",
                    "11.  CDK 인프라 구현 및 검증",
                    "12.  비용 분석: 시나리오 및 가정",
                    "13.  비용 분석: 비교 결과 및 비용 곡선",
                    "14.  결론 및 다음 단계",
                ]

                for i, item in enumerate(toc_items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = item
                    run.font.size = Pt(14)
                    run.font.name = FONT_NAME
                    # Highlight new items
                    if i in (10, 11, 12):  # New sections
                        run.font.color.rgb = BLUE_2
                        run.font.bold = True
                    else:
                        run.font.color.rgb = DARK_TEXT
                break

    # Update page number
    for shape in slide.shapes:
        if shape.has_text_frame and "/" in shape.text_frame.text and shape.text_frame.text.strip().endswith(str(16)):
            shape.text_frame.paragraphs[0].clear()
            run = shape.text_frame.paragraphs[0].add_run()
            run.text = f"2 / {total_slides}"
            run.font.size = Pt(11)
            run.font.color.rgb = GRAY_TEXT
            run.font.name = FONT_NAME
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            break


def update_conclusion_slide(slide, total_slides, page_num):
    """Update conclusion slide to add items 7 and 8."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "11. 결론" in shape.text_frame.text:
                # Update title number
                p = shape.text_frame.paragraphs[0]
                p.clear()
                run = p.add_run()
                run.text = "14. 결론"
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = FONT_NAME

    # Find the content area with numbered items and add new ones
    for shape in slide.shapes:
        if shape.has_text_frame and "하이브리드 프로덕션" in shape.text_frame.text:
            tf = shape.text_frame
            # Add items 6, 7, 8 after existing content
            p6 = tf.add_paragraph()
            run6 = p6.add_run()
            run6.text = ""

            p7h = tf.add_paragraph()
            run7h = p7h.add_run()
            run7h.text = "6   CDK 인프라 구현 및 배포 완료"
            run7h.font.size = Pt(13)
            run7h.font.bold = True
            run7h.font.name = FONT_NAME
            run7h.font.color.rgb = BLUE_2

            p7b = tf.add_paragraph()
            run7b = p7b.add_run()
            run7b.text = "5개 CDK 스택 배포, DynamoDB 시딩 → Lambda EID 해석 → Mountpoint S3\n→ Symlink → samtools CRAM 읽기 End-to-End 검증 완료"
            run7b.font.size = Pt(10)
            run7b.font.name = FONT_NAME
            run7b.font.color.rgb = DARK_TEXT

            p8 = tf.add_paragraph()
            run8 = p8.add_run()
            run8.text = ""

            p8h = tf.add_paragraph()
            run8h = p8h.add_run()
            run8h.text = "7   비용 분석: Mountpoint S3 방식이 96~99.9% 비용 절감"
            run8h.font.size = Pt(13)
            run8h.font.bold = True
            run8h.font.name = FONT_NAME
            run8h.font.color.rgb = GREEN_ACCENT

            p8b = tf.add_paragraph()
            run8b = p8b.add_run()
            run8b.text = "100명 × 15TB: 월 $464 (EBS 복사 $142,303 대비 307배 저렴)\n연간 절감액: EBS 대비 $1.7M, EFS 대비 $6.7M"
            run8b.font.size = Pt(10)
            run8b.font.name = FONT_NAME
            run8b.font.color.rgb = DARK_TEXT
            break

    # Update page number
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "/" in text and text.startswith("1"):
                try:
                    parts = text.split("/")
                    old_num = int(parts[0].strip())
                    if old_num == 14:
                        shape.text_frame.paragraphs[0].clear()
                        run = shape.text_frame.paragraphs[0].add_run()
                        run.text = f"{page_num} / {total_slides}"
                        run.font.size = Pt(11)
                        run.font.color.rgb = GRAY_TEXT
                        run.font.name = FONT_NAME
                        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                except:
                    pass


def update_next_steps_slide(slide, total_slides, page_num):
    """Update 'next steps' slide to mark completed items."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text
                # Mark completed items
                if "비용 분석" in text and "S3 GET" in text:
                    para.clear()
                    run = para.add_run()
                    run.text = "•  비용 분석: S3 GET 요청 수/데이터 전송량 기반 비용 산출 ✓ 완료"
                    run.font.size = Pt(13)
                    run.font.name = FONT_NAME
                    run.font.color.rgb = GREEN_ACCENT
                elif "DynamoDB 연동" in text:
                    para.clear()
                    run = para.add_run()
                    run.text = "•  DynamoDB 연동: 동적 EID 매핑 DB 구현 ✓ 완료 (CDK 배포)"
                    run.font.size = Pt(13)
                    run.font.name = FONT_NAME
                    run.font.color.rgb = GREEN_ACCENT
                elif "Lambda 세션 초기화" in text:
                    para.clear()
                    run = para.add_run()
                    run.text = "•  Lambda 세션 초기화: 연구자 로그인 시 EID symlink 자동 생성 ✓ 완료"
                    run.font.size = Pt(13)
                    run.font.name = FONT_NAME
                    run.font.color.rgb = GREEN_ACCENT
                elif "Cognito 통합" in text:
                    para.clear()
                    run = para.add_run()
                    run.text = "•  Cognito 통합: 외부 연구자 인증 + 프로젝트별 접근 제어 ✓ 완료 (CDK 배포)"
                    run.font.size = Pt(13)
                    run.font.name = FONT_NAME
                    run.font.color.rgb = GREEN_ACCENT
                elif "CloudTrail 감사 로그" in text:
                    para.clear()
                    run = para.add_run()
                    run.text = "•  CloudTrail 감사 로그: CRAM 파일 접근 이력 추적 ✓ 완료 (CDK 배포)"
                    run.font.size = Pt(13)
                    run.font.name = FONT_NAME
                    run.font.color.rgb = GREEN_ACCENT

    # Update page number
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text == "15 / 16":
                shape.text_frame.paragraphs[0].clear()
                run = shape.text_frame.paragraphs[0].add_run()
                run.text = f"{page_num} / {total_slides}"
                run.font.size = Pt(11)
                run.font.color.rgb = GRAY_TEXT
                run.font.name = FONT_NAME
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def update_page_numbers(slides, total_slides):
    """Update page numbers on all existing content slides."""
    for i, slide in enumerate(slides):
        page_num = i + 1  # 1-indexed
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Match pattern "N / 16"
                if " / 16" in text:
                    try:
                        old_num = int(text.split("/")[0].strip())
                        # Don't update slides we handle separately
                        if old_num not in (2, 14, 15):
                            shape.text_frame.paragraphs[0].clear()
                            run = shape.text_frame.paragraphs[0].add_run()
                            run.text = f"{page_num} / {total_slides}"
                            run.font.size = Pt(11)
                            run.font.color.rgb = GRAY_TEXT
                            run.font.name = FONT_NAME
                            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                    except ValueError:
                        pass


def move_slide_to_position(prs, slide, position):
    """Move a slide to a specific position (0-indexed) in the presentation."""
    slide_list = prs.slides._sldIdLst
    # Get the slide's sldId element
    slide_id = None
    for sld_id in slide_list:
        if sld_id.get(qn('r:id')) == slide.part.partname.relative_ref(prs.slides._sldIdLst.getparent().getparent()):
            slide_id = sld_id
            break

    # Simpler approach: manipulate the XML directly
    sld_id_lst = prs.slides._sldIdLst
    # Find the element for our slide (it's the last one since we just added it)
    elements = list(sld_id_lst)
    last_elem = elements[-1]
    sld_id_lst.remove(last_elem)

    if position >= len(elements) - 1:
        sld_id_lst.append(last_elem)
    else:
        sld_id_lst.insert(position, last_elem)


def main():
    SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
    prs = Presentation(os.path.join(SCRIPT_DIR, 'PRESENTATION.pptx'))

    # Total slides will be 20 (16 existing + 4 new)
    TOTAL = 20

    # Get blank layout
    blank_layout = prs.slide_layouts[6]  # Usually blank

    # ── Add 4 new slides ──
    # They'll be added at the end (positions 17-20), then moved

    # Slide A: CDK Infrastructure (will become slide 14, page 14)
    slide_cdk = prs.slides.add_slide(blank_layout)
    build_cdk_infra_slide(slide_cdk, 14, TOTAL)
    # Move to position 13 (after slide 13, 0-indexed)
    move_slide_to_position(prs, slide_cdk, 13)

    # Slide B: Cost Overview (will become slide 15, page 15)
    slide_cost1 = prs.slides.add_slide(blank_layout)
    build_cost_overview_slide(slide_cost1, 15, TOTAL)
    move_slide_to_position(prs, slide_cost1, 14)

    # Slide C: Cost Results (will become slide 16, page 16)
    slide_cost2 = prs.slides.add_slide(blank_layout)
    build_cost_results_slide(slide_cost2, 16, TOTAL)
    move_slide_to_position(prs, slide_cost2, 15)

    # Slide D: Cost Curve (will become slide 17, page 17)
    slide_cost3 = prs.slides.add_slide(blank_layout)
    build_cost_curve_slide(slide_cost3, 17, TOTAL)
    move_slide_to_position(prs, slide_cost3, 16)

    # ── Update existing slides ──
    slides = list(prs.slides)

    # Update TOC (slide index 1)
    update_toc_slide(slides[1], TOTAL)

    # Update conclusion slide (now at index 17, was slide 14/index 13)
    update_conclusion_slide(slides[17], TOTAL, 18)

    # Update next steps slide (now at index 18, was slide 15/index 14)
    update_next_steps_slide(slides[18], TOTAL, 19)

    # Update all page numbers on unchanged slides
    update_page_numbers(slides, TOTAL)

    # ── Save ──
    output_path = os.path.join(SCRIPT_DIR, 'PRESENTATION.pptx')
    prs.save(output_path)
    print(f"Saved updated presentation to {output_path}")
    print(f"Total slides: {len(list(prs.slides))}")


if __name__ == "__main__":
    main()
